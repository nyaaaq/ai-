import os
import logging
from werkzeug.utils import secure_filename
from werkzeug.datastructures import FileStorage
import PyPDF2
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
from typing import Optional
from config import Config

logger = logging.getLogger(__name__)

class FileHandler:
    """文件处理器"""
    
    def __init__(self, upload_folder: str):
        self.upload_folder = upload_folder
        os.makedirs(upload_folder, exist_ok=True)
    
    def allowed_file(self, filename: str) -> bool:
        """检查文件类型是否允许"""
        return '.' in filename and \
               filename.rsplit('.', 1)[1].lower() in Config.ALLOWED_EXTENSIONS
    
    def save_file(self, file: FileStorage) -> str:
        """保存上传的文件"""
        if not self.allowed_file(file.filename):
            raise ValueError(f"不支持的文件类型: {file.filename}")
        
        filename = secure_filename(file.filename)
        # 添加时间戳避免文件名冲突
        import time
        timestamp = str(int(time.time()))
        filename = f"{timestamp}_{filename}"
        
        file_path = os.path.join(self.upload_folder, filename)
        file.save(file_path)
        
        return file_path
    
    def extract_content(self, file_path: str) -> str:
        """从文件中提取内容"""
        _, ext = os.path.splitext(file_path.lower())
        
        try:
            if ext == '.txt':
                return self._extract_from_txt(file_path)
            elif ext == '.pdf':
                return self._extract_from_pdf(file_path)
            elif ext == '.docx':
                return self._extract_from_docx(file_path)
            elif ext == '.pptx':
                return self._extract_from_pptx(file_path)
            elif ext in ['.png', '.jpg', '.jpeg']:
                return self._extract_from_image(file_path)
            else:
                raise ValueError(f"不支持的文件类型: {ext}")
        except Exception as e:
            logger.error(f"提取文件内容失败: {str(e)}")
            raise
    
    def _extract_from_txt(self, file_path: str) -> str:
        """从文本文件提取内容"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """从PDF提取内容"""
        content = []
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                content.append(page.extract_text())
        
        return '\n\n'.join(content)
    
    def _extract_from_docx(self, file_path: str) -> str:
        """从Word文档提取内容"""
        doc = docx.Document(file_path)
        content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content.append(paragraph.text)
        
        return '\n\n'.join(content)
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """从PPT提取内容"""
        prs = Presentation(file_path)
        content = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_content = [f"幻灯片 {slide_num + 1}:"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_content.append(shape.text)
            
            if len(slide_content) > 1:
                content.append('\n'.join(slide_content))
        
        return '\n\n'.join(content)
    
    def _extract_from_image(self, file_path: str) -> str:
        """从图片提取内容（OCR）"""
        try:
            # 打开图片
            image = Image.open(file_path)
            
            # 使用OCR提取文字
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            
            if not text.strip():
                return "图片中未检测到文字内容"
            
            return text
        except Exception as e:
            logger.error(f"OCR识别失败: {str(e)}")
            return "图片文字识别失败，请确保已安装Tesseract OCR"
    
    def cleanup_old_files(self, max_age_hours: int = 24):
        """清理旧文件"""
        import time
        current_time = time.time()
        
        for filename in os.listdir(self.upload_folder):
            file_path = os.path.join(self.upload_folder, filename)
            if os.path.isfile(file_path):
                file_age = current_time - os.path.getmtime(file_path)
                if file_age > max_age_hours * 3600:
                    try:
                        os.remove(file_path)
                        logger.info(f"删除旧文件: {filename}")
                    except Exception as e:
                        logger.error(f"删除文件失败: {str(e)}")